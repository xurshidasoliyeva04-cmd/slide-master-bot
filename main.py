import os, json, sqlite3, logging, asyncio
from datetime import datetime
from io import BytesIO
import google.generativeai as genai

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, CallbackQueryHandler,
    MessageHandler, ConversationHandler, ContextTypes, filters,
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- SOZLAMALAR (Sizning iplaringiz) ---
BOT_TOKEN = "8461901986:AAHIQLMa1RckCqGCU71PJuJZCCnfKdWjYXk"
GEMINI_KEY = "AIzaSyBtUB1yq7lZqF29RPozUiIpj0DT9Rh5eU8"

# Gemini AI ulanishi
genai.configure(api_key=GEMINI_KEY)
ai_model = genai.GenerativeModel('models/gemini-1.5-flash')

# Logging
logging.basicConfig(format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)

# Conversation states
WAITING_NAME, WAITING_PAGES, WAITING_DESIGN, WAITING_TOPIC = range(4)

# Dizaynlar
DESIGNS = {
    1: {"name": "🔵 Ko'k Professional", "primary": "1E3A8A", "secondary": "3B82F6", "accent": "DBEAFE", "text": "1E293B"},
    2: {"name": "🟢 Yashil Tabiat", "primary": "166534", "secondary": "22C55E", "accent": "DCFCE7", "text": "1E293B"},
    3: {"name": "🔴 Qizil Energiya", "primary": "991B1B", "secondary": "EF4444", "accent": "FEE2E2", "text": "1E293B"},
    4: {"name": "🟣 Binafsha Kreativ", "primary": "581C87", "secondary": "A855F7", "accent": "F3E8FF", "text": "1E293B"},
    5: {"name": "🟡 Sariq Quyosh", "primary": "854D0E", "secondary": "EAB308", "accent": "FEF9C3", "text": "1E293B"},
    6: {"name": "⚫ Qora Elegant", "primary": "18181B", "secondary": "3F3F46", "accent": "F4F4F5", "text": "18181B"},
}

# ============ DATABASE ============
def init_db():
    conn = sqlite3.connect("users.db")
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        user_id INTEGER PRIMARY KEY, username TEXT, full_name TEXT,
        credits INTEGER DEFAULT 2, referral_code TEXT UNIQUE,
        referred_by INTEGER, referral_count INTEGER DEFAULT 0, created_at TEXT)""")
    conn.commit()
    conn.close()

def get_user(user_id):
    conn = sqlite3.connect("users.db")
    c = conn.cursor()
    c.execute("SELECT * FROM users WHERE user_id = ?", (user_id,))
    user = c.fetchone()
    conn.close()
    return user

def create_user(user_id, username, full_name, referred_by=None):
    conn = sqlite3.connect("users.db")
    c = conn.cursor()
    referral_code = f"REF{user_id}"
    try:
        c.execute("""INSERT INTO users (user_id, username, full_name, credits, referral_code, referred_by, referral_count, created_at)
            VALUES (?, ?, ?, 2, ?, ?, 0, ?)""", (user_id, username, full_name, referral_code, referred_by, datetime.now().isoformat()))
        if referred_by:
            c.execute("UPDATE users SET referral_count = referral_count + 1, credits = credits + 1 WHERE user_id = ?", (referred_by,))
        conn.commit()
    except sqlite3.IntegrityError: pass
    conn.close()

def use_credit(user_id):
    conn = sqlite3.connect("users.db")
    c = conn.cursor()
    c.execute("UPDATE users SET credits = credits - 1 WHERE user_id = ? AND credits > 0", (user_id,))
    success = c.rowcount > 0
    conn.commit()
    conn.close()
    return success

def get_credits(user_id):
    user = get_user(user_id)
    return user[3] if user else 0

def hex_to_rgb(hex_color):
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

# ============ PPTX GENERATION ============
def create_presentation(topic, full_name, num_pages, design_num, ai_content):
    design = DESIGNS[design_num]
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Bosh slayd
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(design["primary"])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = topic
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Kontent slaydlar (AI matni bilan)
    for i in range(min(len(ai_content), num_pages - 1)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(design["accent"])

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = hex_to_rgb(design["primary"])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8.5), Inches(0.5))
        title_box.text = ai_content[i]['title']
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
        content_box.text = ai_content[i]['content']

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

# ============ HANDLERS ============
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    init_db()
    if not get_user(user.id):
        ref_id = int(context.args[0].replace("REF", "")) if context.args and "REF" in context.args[0] else None
        create_user(user.id, user.username, user.full_name, ref_id)
    
    keyboard = [[InlineKeyboardButton("📊 Slayd Yaratish", callback_data="create_slide")]]
    await update.message.reply_text(f"Xush kelibsiz! Kreditlaringiz: {get_credits(user.id)}", reply_markup=InlineKeyboardMarkup(keyboard))
    return ConversationHandler.END

async def create_slide_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    if get_credits(query.from_user.id) <= 0:
        await query.edit_message_text("Kredit yetarli emas!")
        return ConversationHandler.END
    await query.edit_message_text("Ism-familiyangizni yozing:")
    return WAITING_NAME

async def receive_name(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["full_name"] = update.message.text
    keyboard = [[InlineKeyboardButton(str(i), callback_data=f"pages_{i}") for i in range(8, 12)]]
    await update.message.reply_text("Betlar soni?", reply_markup=InlineKeyboardMarkup(keyboard))
    return WAITING_PAGES

async def receive_pages(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    context.user_data["pages"] = int(query.data.replace("pages_", ""))
    keyboard = [[InlineKeyboardButton(d["name"], callback_data=f"design_{n}")] for n, d in DESIGNS.items()]
    await query.edit_message_text("Dizayn tanlang:", reply_markup=InlineKeyboardMarkup(keyboard))
    return WAITING_DESIGN

async def receive_design(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    context.user_data["design"] = int(query.data.replace("design_", ""))
    await query.edit_message_text("Mavzuni yozing:")
    return WAITING_TOPIC

async def receive_topic(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = update.message.text
    user_id = update.effective_user.id
    msg = await update.message.reply_text("⌛ AI matn tayyorlamoqda...")
    
    try:
        prompt = f"{topic} mavzusida slayd uchun {context.user_data['pages']} ta sarlavha va matn tayyorla. Format: Sarlavha | Matn"
        response = ai_model.generate_content(prompt)
        ai_data = []
        for line in response.text.split('\n'):
            if '|' in line:
                t, c = line.split('|', 1)
                ai_data.append({'title': t.strip(), 'content': c.strip()})

        use_credit(user_id)
        pptx = create_presentation(topic, context.user_data['full_name'], context.user_data['pages'], context.user_data['design'], ai_data)
        await update.message.reply_document(document=BytesIO(pptx.read()), filename=f"{topic}.pptx")
    except Exception as e:
        await update.message.reply_text(f"Xato: {e}")
    return ConversationHandler.END

def main():
    application = Application.builder().token(BOT_TOKEN).build()
    conv_handler = ConversationHandler(
        entry_points=[CallbackQueryHandler(create_slide_start, pattern="create_slide")],
        states={
            WAITING_NAME: [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_name)],
            WAITING_PAGES: [CallbackQueryHandler(receive_pages, pattern="^pages_")],
            WAITING_DESIGN: [CallbackQueryHandler(receive_design, pattern="^design_")],
            WAITING_TOPIC: [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_topic)],
        },
        fallbacks=[CommandHandler("start", start)],
    )
    application.add_handler(CommandHandler("start", start))
    application.add_handler(conv_handler)
    application.run_polling()

if __name__ == "__main__":
    main()
